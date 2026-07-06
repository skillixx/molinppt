from pathlib import Path
import math

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "template-reference-previews" / "high-end-data-viz"
PPTX_PATH = OUT_DIR / "高端数据可视化图表-参考模板原型.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
PNG_W = 1600
PNG_H = 900

BLUE = RGBColor(37, 80, 140)
DEEP_BLUE = RGBColor(29, 62, 110)
MID_BLUE = RGBColor(49, 95, 160)
RED = RGBColor(183, 77, 64)
LIGHT_BG = RGBColor(244, 247, 251)
TEXT = RGBColor(30, 45, 65)
MUTED = RGBColor(103, 116, 137)
WHITE = RGBColor(255, 255, 255)


def font_path(bold=False):
    # 使用微软雅黑，保证中文预览图清晰；如果本机没有该字体，则回退到默认字体。
    candidates = [
        Path("C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def pil_font(size, bold=False):
    selected = font_path(bold)
    if selected:
        return ImageFont.truetype(selected, size)
    return ImageFont.load_default()


def rgb(color):
    return (color[0], color[1], color[2]) if isinstance(color, tuple) else (color.rgb[0], color.rgb[1], color.rgb[2])


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    # PPTX 中的圆角矩形用于数据卡片；普通矩形用于密集看板容器。
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.6)
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide):
    add_rect(slide, 0, 0, SLIDE_W, 1.12, BLUE)
    add_text(slide, "高端数据可视化图表", 0.28, 0.08, 6.8, 0.46, 32, WHITE, True)
    add_text(slide, "述职演讲汇报", 0.32, 0.62, 2.8, 0.28, 17, WHITE, True)
    add_text(slide, "√ 均可编辑", 4.25, 0.62, 2.1, 0.28, 17, WHITE, True)
    add_text(slide, "陇小七", 6.75, 0.62, 1.6, 0.28, 17, WHITE, True)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.42), Inches(0.12), Inches(0.62), Inches(0.62))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(218, 229, 242)
    circle.line.color.rgb = WHITE


def add_dark_blue_background(slide):
    # 首页和结尾页使用独立深蓝商务背景，和中间数据页形成模板层次。
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP_BLUE
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DEEP_BLUE)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(31, 73, 128))
    add_rect(slide, 0, 0, SLIDE_W, 1.25, BLUE)
    add_rect(slide, 0, 6.85, SLIDE_W, 0.65, RGBColor(24, 48, 88))
    for index in range(7):
        x = 8.35 + index * 0.72
        y = 1.35 + (index % 3) * 0.62
        add_rect(slide, x, y, 0.46, 2.8 - index * 0.2, RGBColor(58, 103, 166), line=RGBColor(58, 103, 166))
    for index, size in enumerate([2.25, 1.62, 1.05]):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.4 + index * 0.7), Inches(2.55 + index * 0.46), Inches(size), Inches(size))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(69, 118, 178)
        circle.fill.transparency = 42
        circle.line.color.rgb = RGBColor(110, 154, 205)
    for index in range(5):
        line = slide.shapes.add_connector(
            1,
            Inches(0.65 + index * 1.55),
            Inches(6.12 - index * 0.18),
            Inches(2.0 + index * 1.55),
            Inches(5.42 - index * 0.28),
        )
        line.line.color.rgb = RGBColor(121, 165, 213)
        line.line.width = Pt(1.6)


def create_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_blue_background(slide)
    add_text(slide, "HIGH-END DATA VISUALIZATION", 0.65, 1.46, 5.8, 0.28, 13, RGBColor(196, 216, 239), True)
    add_text(slide, "高端数据\n可视化图表", 0.62, 1.86, 6.2, 1.55, 39, WHITE, True)
    add_text(slide, "述职演讲汇报 · 年度经营分析 · 区域销售复盘", 0.68, 3.55, 5.8, 0.32, 18, RGBColor(225, 235, 247), True)
    add_rect(slide, 0.7, 4.28, 1.45, 0.42, RGBColor(214, 80, 65), radius=True)
    add_text(slide, "均可编辑", 0.92, 4.39, 0.95, 0.14, 12, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, 2.32, 4.28, 1.45, 0.42, RGBColor(70, 129, 190), radius=True)
    add_text(slide, "图表看板", 2.54, 4.39, 0.95, 0.14, 12, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, 3.94, 4.28, 1.45, 0.42, RGBColor(33, 93, 156), radius=True)
    add_text(slide, "商务汇报", 4.16, 4.39, 0.95, 0.14, 12, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, 7.55, 4.88, 4.55, 1.08, RGBColor(244, 247, 251), line=RGBColor(194, 213, 235), radius=True)
    add_text(slide, "2024 DATA REPORT", 7.88, 5.12, 1.9, 0.16, 11, DEEP_BLUE, True)
    add_text(slide, "销售 / 区域 / 产品 / 客户", 7.88, 5.46, 2.25, 0.16, 14, TEXT, True)
    return slide


def create_ending_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_blue_background(slide)
    add_text(slide, "THANKS", 0.78, 1.78, 4.8, 0.55, 45, WHITE, True)
    add_text(slide, "感谢观看", 0.84, 2.72, 3.8, 0.5, 32, WHITE, True)
    add_text(slide, "以数据驱动决策 · 以图表呈现价值", 0.88, 3.48, 5.3, 0.3, 19, RGBColor(226, 237, 249), True)
    add_rect(slide, 0.9, 4.45, 4.3, 0.74, RGBColor(244, 247, 251), line=RGBColor(194, 213, 235), radius=True)
    add_text(slide, "汇报人：陇小七    日期：2024", 1.18, 4.68, 3.4, 0.18, 14, DEEP_BLUE, True)
    add_text(slide, "DATA VISUALIZATION REPORT", 0.9, 6.98, 3.3, 0.15, 10, RGBColor(187, 207, 232), True)
    return slide


def add_section_label(slide, y):
    add_rect(slide, 0.28, y, 2.1, 0.34, DEEP_BLUE)
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(2.16), Inches(y), Inches(0.28), Inches(0.34))
    tri.fill.solid()
    tri.fill.fore_color.rgb = DEEP_BLUE
    tri.line.color.rgb = DEEP_BLUE
    add_text(slide, "数据分析PPT", 0.42, y + 0.06, 1.5, 0.18, 11, WHITE, True)


def add_ring(slide, x, y, value, label, color=MID_BLUE):
    # 用圆环图表达比例，贴近参考图中的 35.28% 与 60.99% 指标样式。
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.28), Inches(1.28))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = RGBColor(213, 222, 235)
    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x + 0.05), Inches(y + 0.05), Inches(1.18), Inches(1.18))
    arc.line.color.rgb = color
    arc.line.width = Pt(5.5)
    add_text(slide, f"{value}%", x + 0.16, y + 0.45, 0.96, 0.25, 18, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.05, y + 1.34, 1.2, 0.18, 8, MUTED, True, PP_ALIGN.CENTER)


def add_bar_chart(slide, x, y, w, h, values, labels, color=DEEP_BLUE):
    max_v = max(values)
    add_text(slide, "各产品销售数量（单位：件）", x, y - 0.28, w, 0.2, 9, TEXT, True)
    for idx, value in enumerate(values):
        bar_w = w / len(values) * 0.48
        gap = w / len(values)
        bx = x + idx * gap + gap * 0.22
        bh = h * value / max_v
        add_rect(slide, bx, y + h - bh, bar_w, bh, color)
        add_text(slide, str(value), bx - 0.03, y + h - bh - 0.18, bar_w + 0.08, 0.12, 6, TEXT, False, PP_ALIGN.CENTER)
        add_text(slide, labels[idx], bx - 0.03, y + h + 0.06, bar_w + 0.08, 0.12, 6, TEXT, False, PP_ALIGN.CENTER)


def add_region_card(slide, x, y, title, target, actual, rate, color):
    add_rect(slide, x, y, 1.05, 2.16, color, radius=True)
    add_text(slide, title, x + 0.1, y + 0.12, 0.85, 0.16, 8, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, f"目标金额\n{target}万", x + 0.16, y + 0.45, 0.75, 0.42, 8, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, f"达成金额\n{actual}万", x + 0.16, y + 0.96, 0.75, 0.42, 8, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, rate, x + 0.16, y + 1.5, 0.75, 0.24, 17, WHITE, True, PP_ALIGN.CENTER)


def create_slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide)

    add_rect(slide, 0.12, 1.34, 13.06, 5.8, WHITE, RGBColor(219, 226, 236))
    add_section_label(slide, 1.38)
    add_rect(slide, 0.38, 1.88, 3.1, 0.32, DEEP_BLUE)
    add_text(slide, "2024年销售基本情况", 0.5, 1.94, 2.2, 0.16, 12, WHITE, True)

    add_ring(slide, 0.7, 2.36, "35.28", "较去年同期提升")
    add_ring(slide, 2.15, 2.36, "60.99", "用户数量增加")

    add_rect(slide, 0.46, 4.25, 3.3, 1.48, WHITE, RGBColor(222, 230, 240))
    add_text(slide, "男女客户占比（单位：%）", 0.56, 4.36, 2.2, 0.18, 9, TEXT, True)
    for i in range(10):
        add_rect(slide, 1.05 + i * 0.22, 4.75, 0.08, 0.38, MID_BLUE if i < 8 else RGBColor(196, 202, 212), radius=True)
        add_rect(slide, 1.05 + i * 0.22, 5.23, 0.08, 0.38, RED if i < 7 else RGBColor(196, 202, 212), radius=True)
    add_text(slide, "男客户", 0.58, 4.78, 0.46, 0.16, 8, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, 0.54, 4.68, 0.5, 0.34, MID_BLUE)
    add_text(slide, "80%", 3.12, 4.84, 0.42, 0.16, 10, MID_BLUE, True)
    add_rect(slide, 0.54, 5.16, 0.5, 0.34, RED)
    add_text(slide, "女客户", 0.58, 5.25, 0.46, 0.16, 8, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "65%", 3.12, 5.32, 0.42, 0.16, 10, RED, True)

    add_rect(slide, 4.05, 1.72, 4.3, 2.15, WHITE, RGBColor(219, 226, 236))
    add_bar_chart(slide, 4.25, 2.15, 3.86, 1.32, [695, 436, 785, 561, 790, 493], ["产品A", "产品B", "产品C", "产品D", "产品E", "产品F"])

    cards = [
        ("陕西", 159, 100, "85%", MID_BLUE),
        ("上海", 190, 180, "96%", RGBColor(55, 108, 170)),
        ("北京", 200, 182, "95%", RED),
        ("深圳", 160, 140, "89%", RGBColor(193, 86, 67)),
    ]
    for idx, card in enumerate(cards):
        add_region_card(slide, 4.16 + idx * 1.2, 4.1, *card)

    return slide


def add_metric_block(slide, x, y, title, value, color):
    add_rect(slide, x, y, 2.85, 1.06, color)
    add_text(slide, title, x + 0.12, y + 0.15, 1.2, 0.15, 8, WHITE, True)
    add_text(slide, f"{value}元", x + 0.12, y + 0.48, 2.1, 0.26, 21, WHITE, True)


def add_combo_chart(slide, x, y):
    add_rect(slide, x, y, 5.52, 2.86, WHITE, RGBColor(219, 226, 236))
    add_text(slide, "销售收入趋势与达成率", x + 0.18, y + 0.14, 1.8, 0.18, 9, TEXT, True)
    values = [5200, 6800, 2100, 10000, 5400, 7600, 12400]
    max_v = max(values)
    points = []
    for idx, value in enumerate(values):
        bx = x + 0.55 + idx * 0.66
        bh = 1.65 * value / max_v
        add_rect(slide, bx, y + 2.25 - bh, 0.22, bh, DEEP_BLUE)
        px = bx + 0.11
        py = y + 2.12 - (0.95 + math.sin(idx * 1.2) * 0.45)
        points.append((px, py))
    for idx in range(len(points) - 1):
        line = slide.shapes.add_connector(1, Inches(points[idx][0]), Inches(points[idx][1]), Inches(points[idx + 1][0]), Inches(points[idx + 1][1]))
        line.line.color.rgb = RED
        line.line.width = Pt(2)
    for px, py in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.04), Inches(py - 0.04), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = RED


def create_slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide)

    add_rect(slide, 0.12, 1.34, 13.06, 5.8, WHITE, RGBColor(219, 226, 236))
    add_section_label(slide, 1.42)

    add_metric_block(slide, 0.35, 1.92, "华北区域", "698,397", MID_BLUE)
    add_metric_block(slide, 3.25, 1.92, "西北区域", "874,638", RED)
    add_metric_block(slide, 6.15, 1.92, "华南区域", "974,574", MID_BLUE)
    add_metric_block(slide, 9.05, 1.92, "东北区域", "1,954,663", RED)

    add_rect(slide, 0.35, 3.1, 2.2, 3.42, WHITE, RGBColor(219, 226, 236))
    add_rect(slide, 0.35, 3.1, 2.2, 0.34, RED)
    add_text(slide, "收入来源", 0.55, 3.2, 0.7, 0.12, 8, WHITE, True)
    add_text(slide, "金额", 1.85, 3.2, 0.36, 0.12, 8, WHITE, True)
    rows = ["产品销售", "服务收入", "会员订阅", "渠道返利", "广告营销", "售后收入", "其他收入"]
    for idx, row in enumerate(rows):
        yy = 3.62 + idx * 0.36
        add_text(slide, row, 0.56, yy, 0.8, 0.12, 7, TEXT)
        add_text(slide, str(53000 - idx * 4300), 1.7, yy, 0.48, 0.12, 7, TEXT, False, PP_ALIGN.RIGHT)

    add_combo_chart(slide, 2.72, 3.1)

    add_rect(slide, 8.45, 3.1, 4.28, 1.42, WHITE, RGBColor(219, 226, 236))
    add_bar_chart(slide, 8.72, 3.48, 3.65, 0.74, [688, 872, 934, 1554], ["一季度", "二季度", "三季度", "四季度"], RED)

    add_rect(slide, 8.45, 4.72, 4.28, 1.8, WHITE, RGBColor(219, 226, 236))
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(9.05), Inches(4.98), Inches(1.65), Inches(1.25))
    pie.fill.solid()
    pie.fill.fore_color.rgb = MID_BLUE
    pie.line.color.rgb = WHITE
    add_text(slide, "区域收入结构", 10.9, 5.0, 1.1, 0.16, 9, TEXT, True)
    for idx, (name, color) in enumerate([("华北", MID_BLUE), ("华南", RED), ("华东", RGBColor(78, 111, 168))]):
        add_rect(slide, 10.9, 5.35 + idx * 0.28, 0.16, 0.12, color)
        add_text(slide, name, 11.12, 5.32 + idx * 0.28, 0.5, 0.12, 7, TEXT)

    return slide


def draw_png_header(draw):
    draw.rectangle([0, 0, PNG_W, 118], fill=(37, 80, 140))
    draw.text((36, 18), "高端数据可视化图表", font=pil_font(58, True), fill=(255, 255, 255))
    draw.text((42, 82), "述职演讲汇报", font=pil_font(28, True), fill=(255, 255, 255))
    draw.text((570, 82), "√ 均可编辑", font=pil_font(28, True), fill=(255, 255, 255))
    draw.text((930, 82), "陇小七", font=pil_font(28, True), fill=(255, 255, 255))


def draw_dark_png_background(draw):
    # PNG 预览中的首页和结尾页背景，与 PPTX 深蓝商务背景保持同一视觉方向。
    draw.rectangle([0, 0, PNG_W, PNG_H], fill=(31, 73, 128))
    draw.rectangle([0, 0, PNG_W, 150], fill=(37, 80, 140))
    draw.rectangle([0, 820, PNG_W, PNG_H], fill=(24, 48, 88))
    for i in range(7):
        x = 980 + i * 82
        y = 168 + (i % 3) * 72
        draw.rounded_rectangle([x, y, x + 52, 690 - i * 16], radius=8, fill=(58, 103, 166))
    for bbox, outline in [
        ([1060, 318, 1320, 578], (110, 154, 205)),
        ([1160, 390, 1348, 578], (121, 165, 213)),
        ([1238, 448, 1362, 572], (146, 184, 222)),
    ]:
        draw.ellipse(bbox, outline=outline, width=10)
    for i in range(5):
        draw.line([80 + i * 184, 730 - i * 20, 235 + i * 184, 650 - i * 32], fill=(121, 165, 213), width=3)


def draw_cover_png():
    image = Image.new("RGB", (PNG_W, PNG_H), (31, 73, 128))
    draw = ImageDraw.Draw(image)
    draw_dark_png_background(draw)
    draw.text((82, 174), "HIGH-END DATA VISUALIZATION", font=pil_font(28, True), fill=(196, 216, 239))
    draw.text((78, 235), "高端数据", font=pil_font(84, True), fill=(255, 255, 255))
    draw.text((78, 335), "可视化图表", font=pil_font(84, True), fill=(255, 255, 255))
    draw.text((84, 462), "述职演讲汇报 · 年度经营分析 · 区域销售复盘", font=pil_font(34, True), fill=(225, 235, 247))
    for i, (label, color) in enumerate([("均可编辑", (214, 80, 65)), ("图表看板", (70, 129, 190)), ("商务汇报", (33, 93, 156))]):
        x = 86 + i * 190
        draw.rounded_rectangle([x, 546, x + 160, 598], radius=18, fill=color)
        draw.text((x + 29, 559), label, font=pil_font(24, True), fill=(255, 255, 255))
    draw.rounded_rectangle([900, 596, 1448, 728], radius=24, fill=(244, 247, 251), outline=(194, 213, 235), width=3)
    draw.text((945, 628), "2024 DATA REPORT", font=pil_font(24, True), fill=(29, 62, 110))
    draw.text((945, 674), "销售 / 区域 / 产品 / 客户", font=pil_font(31, True), fill=(30, 45, 65))
    path = OUT_DIR / "高端数据可视化图表-预览0-首页.png"
    image.save(path)
    return path


def draw_ending_png():
    image = Image.new("RGB", (PNG_W, PNG_H), (31, 73, 128))
    draw = ImageDraw.Draw(image)
    draw_dark_png_background(draw)
    draw.text((96, 218), "THANKS", font=pil_font(92, True), fill=(255, 255, 255))
    draw.text((104, 338), "感谢观看", font=pil_font(62, True), fill=(255, 255, 255))
    draw.text((108, 446), "以数据驱动决策 · 以图表呈现价值", font=pil_font(34, True), fill=(226, 237, 249))
    draw.rounded_rectangle([108, 560, 626, 650], radius=24, fill=(244, 247, 251), outline=(194, 213, 235), width=3)
    draw.text((150, 590), "汇报人：陇小七    日期：2024", font=pil_font(28, True), fill=(29, 62, 110))
    draw.text((108, 838), "DATA VISUALIZATION REPORT", font=pil_font(22, True), fill=(187, 207, 232))
    path = OUT_DIR / "高端数据可视化图表-预览3-结尾.png"
    image.save(path)
    return path


def draw_label(draw, x, y):
    draw.rounded_rectangle([x, y, x + 290, y + 42], radius=2, fill=(29, 62, 110))
    draw.polygon([(x + 290, y), (x + 330, y + 21), (x + 290, y + 42)], fill=(29, 62, 110))
    draw.text((x + 26, y + 8), "数据分析PPT", font=pil_font(20, True), fill=(255, 255, 255))


def draw_chart_mockups(draw, slide_index):
    draw.rounded_rectangle([10, 130, 1590, 885], radius=0, fill=(255, 255, 255), outline=(220, 226, 235), width=3)
    draw_label(draw, 28, 150)
    if slide_index == 1:
        draw.rectangle([48, 220, 510, 262], fill=(29, 62, 110))
        draw.text((64, 227), "2024年销售基本情况", font=pil_font(24, True), fill=(255, 255, 255))
        for cx, value, label in [(165, "35.28%", "较去年同期提升"), (370, "60.99%", "用户数量增加")]:
            draw.ellipse([cx - 78, 306, cx + 78, 462], outline=(49, 95, 160), width=12)
            draw.text((cx - 62, 356), value, font=pil_font(31, True), fill=(29, 62, 110))
            draw.text((cx - 78, 492), label, font=pil_font(18, True), fill=(70, 85, 108))
        draw.rounded_rectangle([48, 560, 510, 780], radius=6, fill=(255, 255, 255), outline=(220, 226, 235), width=2)
        draw.text((68, 582), "男女客户占比（单位：%）", font=pil_font(20, True), fill=(30, 45, 65))
        for row, color, text, pct in [(0, (49, 95, 160), "男客户", "80%"), (1, (183, 77, 64), "女客户", "65%")]:
            yy = 640 + row * 68
            draw.rounded_rectangle([72, yy, 150, yy + 42], radius=2, fill=color)
            draw.text((84, yy + 9), text, font=pil_font(16, True), fill=(255, 255, 255))
            for i in range(10):
                icon_color = color if i < (8 if row == 0 else 7) else (200, 205, 214)
                draw.rectangle([180 + i * 28, yy + 4, 192 + i * 28, yy + 42], fill=icon_color)
            draw.text((452, yy + 8), pct, font=pil_font(22, True), fill=color)
        draw.rounded_rectangle([570, 192, 1180, 460], radius=4, fill=(255, 255, 255), outline=(220, 226, 235), width=2)
        draw.text((595, 210), "各产品销售数量（单位：件）", font=pil_font(18, True), fill=(30, 45, 65))
        for i, v in enumerate([695, 436, 785, 561, 790, 493]):
            h = int(v / 790 * 150)
            x = 620 + i * 86
            draw.rectangle([x, 405 - h, x + 34, 405], fill=(29, 62, 110))
            draw.text((x - 4, 415), f"产品{chr(65 + i)}", font=pil_font(13), fill=(30, 45, 65))
        for i, (city, color, rate) in enumerate([("陕西", (49, 95, 160), "85%"), ("上海", (55, 108, 170), "96%"), ("北京", (183, 77, 64), "95%"), ("深圳", (193, 86, 67), "89%")]):
            x = 592 + i * 230
            draw.rounded_rectangle([x, 510, x + 190, 790], radius=18, fill=color)
            draw.text((x + 68, 532), city, font=pil_font(22, True), fill=(255, 255, 255))
            draw.text((x + 50, 700), rate, font=pil_font(42, True), fill=(255, 255, 255))
    else:
        for i, (title, value, color) in enumerate([("华北区域", "698,397元", (49, 95, 160)), ("西北区域", "874,638元", (183, 77, 64)), ("华南区域", "974,574元", (49, 95, 160)), ("东北区域", "1,954,663元", (183, 77, 64))]):
            x = 44 + i * 380
            draw.rectangle([x, 224, x + 340, 340], fill=color)
            draw.text((x + 20, 244), title, font=pil_font(18, True), fill=(255, 255, 255))
            draw.text((x + 20, 286), value, font=pil_font(36, True), fill=(255, 255, 255))
        draw.rounded_rectangle([42, 390, 346, 825], radius=4, fill=(255, 255, 255), outline=(220, 226, 235), width=2)
        draw.rectangle([42, 390, 346, 434], fill=(183, 77, 64))
        draw.text((70, 402), "收入来源", font=pil_font(18, True), fill=(255, 255, 255))
        for i, row in enumerate(["产品销售", "服务收入", "会员订阅", "渠道返利", "广告营销", "售后收入", "其他收入"]):
            draw.text((72, 462 + i * 46), row, font=pil_font(17), fill=(30, 45, 65))
            draw.text((245, 462 + i * 46), str(53000 - i * 4300), font=pil_font(17), fill=(30, 45, 65))
        draw.rounded_rectangle([380, 390, 1030, 825], radius=4, fill=(255, 255, 255), outline=(220, 226, 235), width=2)
        base_y = 770
        points = []
        for i, v in enumerate([5200, 6800, 2100, 10000, 5400, 7600, 12400]):
            x = 450 + i * 76
            h = int(v / 12400 * 260)
            draw.rectangle([x, base_y - h, x + 30, base_y], fill=(29, 62, 110))
            points.append((x + 15, base_y - h + 48))
        draw.line(points, fill=(183, 77, 64), width=5)
        draw.rounded_rectangle([1070, 390, 1548, 570], radius=4, fill=(255, 255, 255), outline=(220, 226, 235), width=2)
        for i, v in enumerate([688, 872, 934, 1554]):
            x = 1126 + i * 86
            h = int(v / 1554 * 105)
            draw.rectangle([x, 540 - h, x + 42, 540], fill=(183, 77, 64))
        draw.rounded_rectangle([1070, 596, 1548, 825], radius=4, fill=(255, 255, 255), outline=(220, 226, 235), width=2)
        draw.pieslice([1130, 630, 1330, 790], start=0, end=130, fill=(49, 95, 160))
        draw.pieslice([1130, 630, 1330, 790], start=130, end=250, fill=(183, 77, 64))
        draw.pieslice([1130, 630, 1330, 790], start=250, end=360, fill=(78, 111, 168))
        draw.text((1370, 660), "区域收入结构", font=pil_font(22, True), fill=(30, 45, 65))


def create_png(slide_index):
    # PNG 用于给用户快速验收视觉方向，PPTX 文件则保留可编辑形状。
    image = Image.new("RGB", (PNG_W, PNG_H), (244, 247, 251))
    draw = ImageDraw.Draw(image)
    draw_png_header(draw)
    draw_chart_mockups(draw, slide_index)
    path = OUT_DIR / f"高端数据可视化图表-预览{slide_index}-内容页.png"
    image.save(path)
    return path


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    create_cover_slide(prs)
    create_slide_one(prs)
    create_slide_two(prs)
    create_ending_slide(prs)
    prs.save(PPTX_PATH)
    png_paths = [draw_cover_png(), create_png(1), create_png(2), draw_ending_png()]
    print(PPTX_PATH)
    for path in png_paths:
        print(path)


if __name__ == "__main__":
    main()
